#!/usr/bin/env python3
"""
Extract content from PowerPoint (.pptx) files into structured JSON.
Outputs slide text, images, and speaker notes for conversion to HTML presentations.

Usage:
    python extract-pptx.py input.pptx [output_dir]

Dependencies:
    pip install python-pptx
"""

import json
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)


def extract_presentation(pptx_path, output_dir=None):
    """Extract all content from a PowerPoint file."""
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)

    if output_dir is None:
        output_dir = pptx_path.parent / f"{pptx_path.stem}-extracted"
    else:
        output_dir = Path(output_dir)

    assets_dir = output_dir / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(pptx_path))
    slides_data = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_data = {
            "slide_number": slide_idx,
            "title": None,
            "body": [],
            "images": [],
            "notes": None,
            "layout": slide.slide_layout.name if slide.slide_layout else None,
        }

        # Extract shapes
        for shape in slide.shapes:
            # Title
            if shape.has_text_frame and shape.shape_id == slide.shapes.title_shape_id if hasattr(slide.shapes, 'title_shape_id') else False:
                slide_data["title"] = shape.text_frame.text.strip()
            elif hasattr(shape, "is_placeholder") and shape.is_placeholder:
                ph = shape.placeholder_format
                if ph.idx == 0:  # Title placeholder
                    if shape.has_text_frame:
                        slide_data["title"] = shape.text_frame.text.strip()
                    continue

            # Text content
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        # Check if this is the title (first text if no title found)
                        if slide_data["title"] is None and paragraph.level == 0:
                            slide_data["title"] = text
                        else:
                            slide_data["body"].append({
                                "text": text,
                                "level": paragraph.level,
                                "bold": any(run.font.bold for run in paragraph.runs if run.font.bold),
                            })

            # Images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                img_filename = f"slide{slide_idx}_img{len(slide_data['images']) + 1}.{image.content_type.split('/')[-1]}"
                img_path = assets_dir / img_filename

                with open(img_path, "wb") as f:
                    f.write(image.blob)

                slide_data["images"].append({
                    "filename": img_filename,
                    "path": str(img_path),
                    "width": shape.width,
                    "height": shape.height,
                    "content_type": image.content_type,
                })

            # Tables
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                slide_data["body"].append({
                    "type": "table",
                    "data": table_data,
                })

        # Speaker notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_data["notes"] = notes

        slides_data.append(slide_data)

    # Write JSON output
    output_file = output_dir / "extracted-slides.json"
    with open(output_file, "w", encoding="utf-8") as f:
        json.dump(slides_data, f, indent=2, ensure_ascii=False, default=str)

    # Print summary
    print(f"\nExtracted {len(slides_data)} slides from: {pptx_path.name}")
    print(f"Output directory: {output_dir}")
    print(f"JSON file: {output_file}")
    print(f"\nSlide Summary:")
    for slide in slides_data:
        img_count = len(slide["images"])
        title = slide["title"] or "(no title)"
        print(f"  Slide {slide['slide_number']}: {title[:50]}{'...' if len(title) > 50 else ''}"
              f" | {len(slide['body'])} text blocks | {img_count} images")

    return slides_data


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract-pptx.py <input.pptx> [output_dir]")
        sys.exit(1)

    input_file = sys.argv[1]
    output_directory = sys.argv[2] if len(sys.argv) > 2 else None
    extract_presentation(input_file, output_directory)
